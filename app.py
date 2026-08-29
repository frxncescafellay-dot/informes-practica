import os
import io
import json
import zipfile
import threading
import time
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email.mime.text import MIMEText
from email import encoders
from datetime import datetime
import pytz
import pandas as pd
import streamlit as st
import streamlit.components.v1 as components
from PIL import Image
from groq import Groq
import pypdf
from docx import Document
from pptx import Presentation

# --- CONFIGURACION DE PAGINA ---
st.set_page_config(
    page_title="Plataforma de Analitica & Materiales IA",
    page_icon="⚡",
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- RUTAS DE ALMACENAMIENTO PERSISTENTE ---
DIR_BASE = "almacen_datos"
DIR_DATASETS = os.path.join(DIR_BASE, "datasets")
DIR_MATERIAL = os.path.join(DIR_BASE, "material")
DIR_INFORMES = os.path.join(DIR_BASE, "informes_guardados")
DIR_USUARIOS = os.path.join(DIR_BASE, "usuarios_guardados")
DIR_AVATARS = os.path.join(DIR_BASE, "avatares")
FILE_DB = os.path.join(DIR_BASE, "base_datos.json")

for d in [DIR_BASE, DIR_DATASETS, DIR_MATERIAL, DIR_INFORMES, DIR_USUARIOS, DIR_AVATARS]:
    os.makedirs(d, exist_ok=True)[cite: 16]

MODELO_WHISPER = "whisper-large-v3"[cite: 16]
CORREO_DESTINO_BACKUP = "francesca.fellay.b@mail.pucv.cl"[cite: 16]

# --- FUNCIONES DE RESPALDO Y ZIP COMPLETO ---
def generar_zip_respaldo_completo():
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
        if os.path.exists(DIR_BASE):
            for root, _, files in os.walk(DIR_BASE):
                for file in files:
                    file_path = os.path.join(root, file)
                    arcname = os.path.relpath(file_path, start=DIR_BASE)
                    zip_file.write(file_path, arcname=arcname)
    zip_buffer.seek(0)
    return zip_buffer.getvalue()[cite: 16]

def restaurar_desde_zip_completo(archivo_zip_bytes):
    with zipfile.ZipFile(io.BytesIO(archivo_zip_bytes), "r") as zip_file:
        zip_file.extractall(DIR_BASE)[cite: 16]

def enviar_correo_backup(zip_bytes, fecha_str):
    smtp_server = st.secrets.get("SMTP_SERVER", os.environ.get("SMTP_SERVER", "smtp.gmail.com"))
    smtp_port = int(st.secrets.get("SMTP_PORT", os.environ.get("SMTP_PORT", 587)))
    smtp_user = st.secrets.get("SMTP_USER", os.environ.get("SMTP_USER", ""))
    smtp_password = st.secrets.get("SMTP_PASSWORD", os.environ.get("SMTP_PASSWORD", ""))
    
    if not smtp_user or not smtp_password:
        return False, "Faltan credenciales SMTP (SMTP_USER y SMTP_PASSWORD en Secrets)"[cite: 16]

    try:
        msg = MIMEMultipart()
        msg["From"] = smtp_user
        msg["To"] = CORREO_DESTINO_BACKUP
        msg["Subject"] = f"📦 Backup Diario Automatizado — Plataforma ({fecha_str})"

        cuerpo = f"""
        Estimada Francesca Fellay,

        Adjunto se encuentra el respaldo diario completo de la base de datos y archivos físicos 
        (Datasets, Materiales, Informes, Usuarios y Avatares) correspondiente a las 08:00 AM ({fecha_str}).

        Este archivo .zip único puede ser descargado y cargado en el sistema en cualquier momento 
        para restaurar todos los registros si la plataforma se reinicia.
        """
        msg.attach(MIMEText(cuerpo, "plain", "utf-8"))

        part = MIMEBase("application", "octet-stream")
        part.set_payload(zip_bytes)
        encoders.encode_base64(part)
        part.add_header("Content-Disposition", f"attachment; filename=backup_plataforma_{datetime.now().strftime('%Y%m%d')}.zip")
        msg.attach(part)

        server = smtplib.SMTP(smtp_server, smtp_port)
        server.starttls()
        server.login(smtp_user, smtp_password)
        server.send_message(msg)
        server.quit()
        return True, "Enviado con éxito"[cite: 16]
    except Exception as e:
        return False, str(e)[cite: 16]

# --- TAREA EN SEGUNDO PLANO: ENVIO AUTOMATICO A LAS 08:00 AM (HORA CHILE) ---
if "scheduler_iniciado" not in st.session_state:
    st.session_state.scheduler_iniciado = True[cite: 16]
    def daemon_backup_diario():
        ultimo_dia_enviado = None
        tz_chile = pytz.timezone("America/Santiago")
        while True:
            ahora = datetime.now(tz_chile)
            dia_actual = ahora.strftime("%Y-%m-%d")
            if ahora.hour == 8 and ultimo_dia_enviado != dia_actual:
                zip_data = generar_zip_respaldo_completo()
                exito, _ = enviar_correo_backup(zip_data, ahora.strftime("%d/%m/%Y %H:%M"))
                if exito:
                    ultimo_dia_enviado = dia_actual
            time.sleep(30)
            
    hilo = threading.Thread(target=daemon_backup_diario, daemon=True)
    hilo.start()[cite: 16]

# --- PALETA TEAL + NARANJA TERRACOTA ---
st.markdown("""
<style>
    html, body, [class*="css"], .stApp {
        font-family: 'Segoe UI', system-ui, -apple-system, sans-serif !important;
    }

    .stApp {
        background: #73b5ae !important;
        color: #061e1b;
    }
    
    h1, h2, h3, h4, h5, h6 {
        color: #061e1b !important;
        font-weight: 800;
        letter-spacing: -0.5px;
        font-family: 'Segoe UI', system-ui, -apple-system, sans-serif !important;
    }
    
    section[data-testid="stSidebar"] {
        background-color: #63a59e !important;
        border-right: 2px solid #52948d !important;
    }
    section[data-testid="stSidebar"] * {
        color: #061e1b !important;
        font-weight: 600;
    }

    .modern-card {
        background: #89c7c0 !important;
        border: 2px solid #63a59e !important;
        border-radius: 14px;
        padding: 22px;
        margin-bottom: 20px;
        box-shadow: 0 4px 12px rgba(13, 44, 41, 0.12);
    }

    div[data-testid="stExpander"] {
        background-color: #89c7c0 !important;
        border: 2px solid #63a59e !important;
        border-radius: 10px !important;
    }
    div[data-testid="stExpander"] summary {
        background-color: #7ab8b1 !important;
        color: #061e1b !important;
        border-radius: 8px !important;
        font-weight: 750 !important;
    }
    div[data-testid="stExpander"] summary * {
        color: #061e1b !important;
        font-weight: 750 !important;
    }
    div[data-testid="stExpander"] div[role="region"] {
        background-color: #89c7c0 !important;
        color: #061e1b !important;
    }

    input[type="text"], 
    input[type="password"], 
    textarea {
        background-color: #a2d2cc !important;
        color: #061e1b !important;
        border: 2px solid #52948d !important;
        border-radius: 8px !important;
        font-size: 0.95rem !important;
        font-weight: 600 !important;
    }
    input[type="text"]:focus, 
    input[type="password"]:focus, 
    textarea:focus {
        background-color: #bfe3de !important;
        color: #000000 !important;
        border-color: #c84b1e !important;
        box-shadow: 0 0 0 3px rgba(200, 75, 30, 0.25) !important;
    }

    div[data-baseweb="input"] button,
    div[data-baseweb="input"] div {
        background-color: #a2d2cc !important;
        border: none !important;
        color: #061e1b !important;
    }
    div[data-baseweb="input"] svg {
        fill: #061e1b !important;
        color: #061e1b !important;
    }

    div[data-baseweb="select"] > div {
        background-color: #a2d2cc !important;
        color: #061e1b !important;
        border: 2px solid #52948d !important;
        border-radius: 8px !important;
    }
    div[data-baseweb="select"] * {
        background-color: #a2d2cc !important;
        color: #061e1b !important;
        font-weight: 600 !important;
    }
    div[data-baseweb="select"] svg {
        fill: #061e1b !important;
        color: #061e1b !important;
    }

    div[data-testid="stCheckbox"] span[role="checkbox"] {
        background-color: #a2d2cc !important;
        border: 2px solid #52948d !important;
        border-radius: 5px !important;
    }
    div[data-testid="stCheckbox"] span[role="checkbox"][aria-checked="true"] {
        background-color: #c84b1e !important;
        border-color: #9e3610 !important;
    }
    div[data-testid="stCheckbox"] svg {
        color: #ffffff !important;
        fill: #ffffff !important;
    }
    div[data-testid="stCheckbox"] label span {
        color: #061e1b !important;
        font-weight: 600 !important;
    }

    div[data-testid="stFileUploader"] {
        background-color: #89c7c0 !important;
        border: 2px dashed #c84b1e !important;
        border-radius: 12px !important;
        padding: 14px !important;
    }
    div[data-testid="stFileUploader"] section {
        background-color: #89c7c0 !important;
    }
    div[data-testid="stFileUploader"] section * {
        color: #061e1b !important;
    }
    div[data-testid="stFileUploader"] button {
        background-color: #c84b1e !important;
        border: 2px solid #9e3610 !important;
        border-radius: 8px !important;
        font-weight: 750 !important;
    }
    div[data-testid="stFileUploader"] button p {
        color: #ffffff !important;
    }

    div[data-testid="stDownloadButton"]>button {
        background: #a2d2cc !important;
        color: #061e1b !important;
        border: 2px solid #52948d !important;
        border-radius: 9px !important;
        padding: 9px 20px !important;
        font-weight: 800 !important;
        font-size: 0.95rem !important;
        box-shadow: 0 3px 8px rgba(13, 44, 41, 0.15) !important;
        transition: all 0.2s ease !important;
    }
    div[data-testid="stDownloadButton"]>button:hover {
        background: #bfe3de !important;
        color: #000000 !important;
        border-color: #061e1b !important;
        transform: translateY(-1px);
    }
    div[data-testid="stDownloadButton"]>button p {
        color: #061e1b !important;
        font-weight: 800 !important;
    }

    .stButton>button, 
    div[data-testid="stFormSubmitButton"]>button {
        background: linear-gradient(135deg, #c84b1e 0%, #b23b14 100%) !important;
        color: #ffffff !important;
        border: 2px solid #9e3610 !important;
        border-radius: 9px !important;
        padding: 9px 22px !important;
        font-weight: 800 !important;
        font-size: 0.95rem !important;
        box-shadow: 0 4px 12px rgba(178, 59, 20, 0.35) !important;
        transition: all 0.2s ease !important;
    }
    .stButton>button:hover, 
    div[data-testid="stFormSubmitButton"]>button:hover {
        background: linear-gradient(135deg, #db5928 0%, #c84b1e 100%) !important;
        border-color: #db5928 !important;
        box-shadow: 0 6px 16px rgba(219, 89, 40, 0.45) !important;
        transform: translateY(-1px);
    }
    .stButton>button p,
    div[data-testid="stFormSubmitButton"]>button p {
        color: #ffffff !important;
        font-weight: 800 !important;
    }

    div[data-testid="stDataFrame"], 
    div[data-testid="stDataEditor"],
    div[data-testid="stDataFrame"] > div,
    div[data-testid="stDataEditor"] > div {
        background-color: #89c7c0 !important;
        border: 2px solid #52948d !important;
        border-radius: 10px !important;
    }
    
    .glide-data-grid,
    .gdg-container,
    div[data-testid="stDataFrame"] canvas,
    div[data-testid="stDataEditor"] canvas {
        filter: invert(0.85) sepia(0.3) saturate(2.5) hue-rotate(130deg) brightness(1.05) !important;
        border-radius: 8px !important;
    }

    table {
        width: 100% !important;
        border-collapse: collapse !important;
        margin: 16px 0 !important;
        background-color: #89c7c0 !important;
        border-radius: 8px !important;
        overflow: hidden !important;
        border: 2px solid #52948d !important;
    }
    th {
        background-color: #63a59e !important;
        color: #061e1b !important;
        font-weight: 800 !important;
        padding: 12px 14px !important;
        border: 1px solid #52948d !important;
        text-align: left !important;
    }
    td {
        padding: 10px 14px !important;
        border: 1px solid #52948d !important;
        color: #061e1b !important;
        font-size: 0.95rem !important;
    }
    tr:nth-child(even) {
        background-color: #7ab8b1 !important;
    }

    .highlight-tag {
        background: #a2d2cc;
        border: 1.5px solid #52948d;
        padding: 3px 10px;
        border-radius: 6px;
        font-weight: 700;
        color: #061e1b;
        font-size: 0.92rem;
        display: inline-block;
        font-family: 'Segoe UI', system-ui, -apple-system, sans-serif !important;
    }

    code, pre {
        background-color: #a2d2cc !important;
        color: #061e1b !important;
        border: 1.5px solid #52948d !important;
        border-radius: 6px !important;
        padding: 3px 8px !important;
        font-family: 'Segoe UI', system-ui, -apple-system, sans-serif !important;
        font-size: 0.92rem !important;
        font-weight: 700 !important;
    }
    
    div[data-testid="stBaseButton-secondary"] button {
        background: #e29b9b !important;
        border: 2px solid #c96b6b !important;
        box-shadow: none !important;
    }
    div[data-testid="stBaseButton-secondary"] button:hover {
        background: #d67a7a !important;
    }
    div[data-testid="stBaseButton-secondary"] button p {
        color: #5c0f0f !important;
    }

    button[data-baseweb="tab"] {
        background-color: #63a59e !important;
        color: #061e1b !important;
        border-radius: 8px 8px 0 0 !important;
        font-weight: 750 !important;
        margin-right: 4px !important;
        padding: 9px 18px !important;
        border: 1px solid #52948d !important;
    }
    button[data-baseweb="tab"][aria-selected="true"] {
        background-color: #89c7c0 !important;
        color: #061e1b !important;
        border-bottom: 3.5px solid #c84b1e !important;
    }

    label, p, span {
        color: #061e1b;
        font-weight: 600;
    }

    .badge-role {
        display: inline-block;
        background: #c84b1e;
        color: #ffffff !important;
        padding: 4px 12px;
        border-radius: 14px;
        font-size: 0.82rem;
        font-weight: 800;
        border: 1.5px solid #9e3610;
    }

    .login-container-teal {
        background: #89c7c0 !important;
        border: 2px solid #52948d !important;
        border-radius: 18px;
        padding: 36px 32px;
        box-shadow: 0 10px 30px rgba(13, 44, 41, 0.2);
        margin-top: 40px;
    }
</style>
""", unsafe_allow_html=True)[cite: 16]

# --- RELOJ DUAL ---
def renderizar_reloj_chile():
    html_reloj = """
    <div style="background: linear-gradient(135deg, #c84b1e 0%, #b23b14 100%); border: 2px solid #9e3610; border-radius: 16px; padding: 12px 20px; display: flex; align-items: center; justify-content: space-between; gap: 16px; box-shadow: 0 6px 18px rgba(178, 59, 20, 0.35); max-width: 480px; margin-left: auto; margin-bottom: 15px; font-family: 'Segoe UI', system-ui, sans-serif;">
        <div style="position: relative; width: 60px; height: 60px;">
            <svg id="analog-clock" width="60" height="60" viewBox="0 0 100 100">
                <circle cx="50" cy="50" r="46" fill="#fbe8e1" stroke="#782307" stroke-width="4"/>
                <line x1="50" y1="10" x2="50" y2="16" stroke="#4a1503" stroke-width="3" stroke-linecap="round"/>
                <line x1="90" y1="50" x2="84" y2="50" stroke="#4a1503" stroke-width="3" stroke-linecap="round"/>
                <line x1="50" y1="90" x2="50" y2="84" stroke="#4a1503" stroke-width="3" stroke-linecap="round"/>
                <line x1="10" y1="50" x2="16" y2="50" stroke="#4a1503" stroke-width="3" stroke-linecap="round"/>
                <line id="hour-hand" x1="50" y1="50" x2="50" y2="28" stroke="#260900" stroke-width="5" stroke-linecap="round"/>
                <line id="min-hand" x1="50" y1="50" x2="50" y2="18" stroke="#591905" stroke-width="3.5" stroke-linecap="round"/>
                <line id="sec-hand" x1="50" y1="50" x2="50" y2="14" stroke="#b23b14" stroke-width="2" stroke-linecap="round"/>
                <circle cx="50" cy="50" r="4" fill="#260900"/>
            </svg>
        </div>
        <div style="display: flex; flex-direction: column; justify-content: center;">
            <div style="font-size: 0.76rem; font-weight: 800; color: #fbe8e1; text-transform: uppercase; letter-spacing: 0.8px;">🇨🇱 Hora Oficial de Chile</div>
            <div id="digital-clock" style="font-size: 1.65rem; font-weight: 900; color: #ffffff; line-height: 1.1; text-shadow: 0 2px 4px rgba(0,0,0,0.2);">--:--:--</div>
            <div id="digital-date" style="font-size: 0.82rem; font-weight: 700; color: #fbe8e1; margin-top: 2px;">Cargando fecha...</div>
        </div>
    </div>
    <script>
        function actualizarReloj() {
            const opciones = { timeZone: 'America/Santiago', hour12: false, hour: '2-digit', minute: '2-digit', second: '2-digit' };
            const opcionesFecha = { timeZone: 'America/Santiago', weekday: 'long', year: 'numeric', month: 'long', day: 'numeric' };
            const ahora = new Date();
            const formateadorHora = new Intl.DateTimeFormat('es-CL', opciones);
            const formateadorFecha = new Intl.DateTimeFormat('es-CL', opcionesFecha);
            const partesHora = formateadorHora.formatToParts(ahora);
            let h = 0, m = 0, s = 0;
            partesHora.forEach(p => {
                if (p.type === 'hour') h = parseInt(p.value);
                if (p.type === 'minute') m = parseInt(p.value);
                if (p.type === 'second') s = parseInt(p.value);
            });
            const hStr = h.toString().padStart(2, '0');
            const mStr = m.toString().padStart(2, '0');
            const sStr = s.toString().padStart(2, '0');
            document.getElementById('digital-clock').textContent = `${hStr}:${mStr}:${sStr}`;
            let fechaStr = formateadorFecha.format(ahora);
            fechaStr = fechaStr.charAt(0).toUpperCase() + fechaStr.slice(1);
            document.getElementById('digital-date').textContent = fechaStr;
            const secDeg = s * 6;
            const minDeg = m * 6 + s * 0.1;
            const hourDeg = (h % 12) * 30 + m * 0.5;
            function rotar(elemId, deg) {
                const el = document.getElementById(elemId);
                if (el) el.setAttribute('transform', `rotate(${deg} 50 50)`);
            }
            rotar('sec-hand', secDeg);
            rotar('min-hand', minDeg);
            rotar('hour-hand', hourDeg);
        }
        setInterval(actualizarReloj, 1000);
        actualizarReloj();
    </script>
    """
    components.html(html_reloj, height=105)[cite: 16]

# --- GESTOR DE PERSISTENCIA Y RESTAURACION TOTAL (ANTI-HIBERNACION) ---
def sincronizar_usuarios_fisicos(usuarios_dict):
    for u_k, u_data in usuarios_dict.items():
        ruta_u_json = os.path.join(DIR_USUARIOS, f"{u_k}.json")
        with open(ruta_u_json, "w", encoding="utf-8") as fu:
            json.dump(u_data, fu, ensure_ascii=False, indent=2)[cite: 16]

def restaurar_repositorio_local(data):
    # 1. Recuperar y Sincronizar Usuarios Creados
    if os.path.exists(DIR_USUARIOS):
        for arch in os.listdir(DIR_USUARIOS):
            if arch.endswith(".json"):
                u_login = arch.rsplit(".json", 1)[0]
                ruta_uj = os.path.join(DIR_USUARIOS, arch)
                try:
                    with open(ruta_uj, "r", encoding="utf-8") as fuj:
                        u_info = json.load(fuj)
                        if u_login not in data["usuarios"]:
                            data["usuarios"][u_login] = u_info
                except Exception:
                    pass

    # 2. Recuperar Datasets
    if os.path.exists(DIR_DATASETS):
        for arch in os.listdir(DIR_DATASETS):
            ruta_f = os.path.join(DIR_DATASETS, arch)
            if os.path.isfile(ruta_f):
                partes = arch.split("_", 1)
                titulo_ds = partes[1].rsplit(".", 1)[0] if len(partes) > 1 else arch
                if titulo_ds not in data["datasets"]:
                    data["datasets"][titulo_ds] = {
                        "titulo": titulo_ds,
                        "autor": "Archivo Restaurado",
                        "fecha": datetime.now().strftime("%Y-%m-%d %H:%M"),
                        "ruta": ruta_f
                    }[cite: 16]

    # 3. Recuperar Materiales
    rutas_mat = {item.get("ruta") for item in data.get("material", []) if item.get("ruta")}
    dir_legacy = os.path.join(DIR_BASE, "intervencion")
    if os.path.exists(dir_legacy):
        for f in os.listdir(dir_legacy):
            origen = os.path.join(dir_legacy, f)
            destino = os.path.join(DIR_MATERIAL, f)
            if not os.path.exists(destino):
                os.rename(origen, destino)[cite: 16]
                
    if os.path.exists(DIR_MATERIAL):
        for arch in os.listdir(DIR_MATERIAL):
            ruta_f = os.path.join(DIR_MATERIAL, arch)
            if ruta_f not in rutas_mat and os.path.isfile(ruta_f):
                partes = arch.split("_", 1)
                nombre_orig = partes[1] if len(partes) > 1 else arch
                ext = arch.split(".")[-1].lower()
                data["material"].append({
                    "titulo": f"Material Restaurado ({nombre_orig})",
                    "nombre_original": nombre_orig,
                    "tipo": ext,
                    "autor": "Sistema / Restauración",
                    "fecha": datetime.now().strftime("%Y-%m-%d %H:%M"),
                    "ruta": ruta_f,
                    "resumen": "Archivo físico recuperado automáticamente del almacenamiento persistente."
                })[cite: 16]

    # 4. Recuperar Informes
    if os.path.exists(DIR_INFORMES):
        titulos_inf = {inf.get("titulo") for inf in data.get("informes", [])}
        for arch in os.listdir(DIR_INFORMES):
            if arch.endswith(".json"):
                ruta_inf = os.path.join(DIR_INFORMES, arch)
                try:
                    with open(ruta_inf, "r", encoding="utf-8") as fi:
                        inf_obj = json.load(fi)
                        if inf_obj.get("titulo") not in titulos_inf:
                            data["informes"].append(inf_obj)
                except Exception:
                    pass
    return data[cite: 16]

def cargar_estado():
    # Estructura base de usuarios predeterminados requeridos
    data_defecto = {
        "usuarios": {
            "Francesca Fellay": {
                "nombre": "Francesca Fellay",
                "rol": "Admin",
                "pin": "1953",
                "permiso_editar": True,
                "permiso_eliminar": True,
                "avatar": ""
            },
            "admin 2": {
                "nombre": "Karol Medina",
                "rol": "Usuario",
                "pin": "5678",
                "permiso_editar": False,
                "permiso_eliminar": False,
                "avatar": ""
            },
            "admin 3": {
                "nombre": "Lukas Núñez",
                "rol": "Usuario",
                "pin": "9999",
                "permiso_editar": False,
                "permiso_eliminar": False,
                "avatar": ""
            },
            "admin 4": {
                "nombre": "Francisca Guerra",
                "rol": "Usuario",
                "pin": "0000",
                "permiso_editar": False,
                "permiso_eliminar": False,
                "avatar": ""
            },
            "gerente": {
                "nombre": "Gerente",
                "rol": "Admin",
                "pin": "gerente",
                "permiso_editar": True,
                "permiso_eliminar": True,
                "avatar": ""
            }
        },
        "datasets": {},
        "material": [],
        "informes": []
    }
    
    if not os.path.exists(FILE_DB):
        data = restaurar_repositorio_local(data_defecto)
        sincronizar_usuarios_fisicos(data["usuarios"])
        guardar_estado(data)
        return data[cite: 16]
        
    with open(FILE_DB, "r", encoding="utf-8") as f:
        try:
            data = json.load(f)
        except Exception:
            data = data_defecto[cite: 16]

    if "material" not in data:
        data["material"] = data.pop("intervencion", [])[cite: 16]
    if "datasets" not in data:
        data["datasets"] = {}[cite: 16]
    if "informes" not in data:
        data["informes"] = [][cite: 16]
    if "usuarios" not in data:
        data["usuarios"] = {}[cite: 16]

    # Asegurar usuarios requeridos y actualizar datos base
    for u_k, u_d in data_defecto["usuarios"].items():
        if u_k not in data["usuarios"]:
            data["usuarios"][u_k] = u_d
        else:
            data["usuarios"][u_k]["pin"] = u_d["pin"]
            data["usuarios"][u_k]["rol"] = u_d["rol"]
            data["usuarios"][u_k]["nombre"] = u_d["nombre"]
            if u_k in ["Francesca Fellay", "gerente"]:
                data["usuarios"][u_k]["permiso_editar"] = True
                data["usuarios"][u_k]["permiso_eliminar"] = True
    
    data = restaurar_repositorio_local(data)
    sincronizar_usuarios_fisicos(data["usuarios"])
    guardar_estado(data)
    return data[cite: 16]

def guardar_estado(data):
    sincronizar_usuarios_fisicos(data.get("usuarios", {}))
    with open(FILE_DB, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)[cite: 16]

db = cargar_estado()[cite: 16]

# --- CLIENTE GROQ IA ---
def obtener_api_key_groq():
    if "GROQ_API_KEY" in st.secrets:
        return st.secrets["GROQ_API_KEY"][cite: 16]
    return os.environ.get("GROQ_API_KEY", "")[cite: 16]

API_KEY_GROQ = obtener_api_key_groq()[cite: 16]

def obtener_cliente_ia():
    if not API_KEY_GROQ:
        return None
    try:
        return Groq(api_key=API_KEY_GROQ)[cite: 16]
    except Exception:
        return None

def ejecutar_chat_groq(client, prompt_sistema, prompt_usuario):
    modelos_candidatos = [
        "llama-3.3-70b-versatile",
        "llama-3.1-70b-versatile",
        "llama3-70b-8192",
        "mixtral-8x7b-32768",
        "gemma2-9b-it"
    ][cite: 16]
    try:
        lista_api = [m.id for m in client.models.list().data if "whisper" not in m.id.lower() and "guard" not in m.id.lower()][cite: 16]
    except Exception:
        lista_api = [][cite: 16]

    candidatos = [m for m in modelos_candidatos if m in lista_api] + lista_api + modelos_candidatos[cite: 16]
    candidatos = list(dict.fromkeys(candidatos))[cite: 16]

    ultimo_error = None
    for model_id in candidatos:
        try:
            resp = client.chat.completions.create(
                model=model_id,
                messages=[
                    {"role": "system", "content": prompt_sistema},
                    {"role": "user", "content": prompt_usuario}
                ],
                temperature=0.3
            )[cite: 16]
            return resp.choices[0].message.content[cite: 16]
        except Exception as e:
            ultimo_error = e
            continue
    raise Exception(f"No fue posible conectar con los modelos de Groq. Detalle: {ultimo_error}")[cite: 16]

# --- EXTRACCION DE TEXTO ---
def extraer_texto_archivo(ruta, extension):
    texto = ""
    try:
        if extension == "pdf":
            lector = pypdf.PdfReader(ruta)
            for pag in lector.pages:
                t = pag.extract_text()
                if t:
                    texto += t + "\n"
        elif extension == "docx":
            doc = Document(ruta)
            texto = "\n".join([p.text for p in doc.paragraphs])
        elif extension == "pptx":
            prs = Presentation(ruta)
            for diap in prs.slides:
                for forma in diap.shapes:
                    if hasattr(forma, "text"):
                        texto += forma.text + "\n"
        elif extension in ["xlsx", "xls"]:
            df_tmp = pd.read_excel(ruta)
            texto = f"Estadisticas:\n{df_tmp.describe(include='all').to_string()}\nPrimeras filas:\n{df_tmp.head(10).to_string()}"
    except Exception as e:
        texto = f"Error al extraer texto: {e}"
    return texto[:8000][cite: 16]

# --- CONTROL DE ACCESO (LOGIN) ---
if "autenticado" not in st.session_state:
    st.session_state.autenticado = False
    st.session_state.usuario_clave = None[cite: 16]

if not st.session_state.autenticado:
    c1, c2, c3 = st.columns([1, 1.4, 1])[cite: 16]
    with c2:
        st.markdown("""
        <div class='login-container-teal'>
            <div style='text-align: center; margin-bottom: 22px;'>
                <div style='font-size: 3rem; margin-bottom: 4px;'>⚡</div>
                <h2 style='margin: 0; color: #061e1b !important; font-size: 1.8rem;'>Acceso a la Plataforma</h2>
                <p style='margin: 4px 0 0 0; color: #061e1b; font-size: 0.95rem; font-weight: 700;'>Gestión & Analítica con Groq IA</p>
            </div>
        """, unsafe_allow_html=True)[cite: 16]
        
        with st.form("login_form"):
            st.markdown("<p style='color:#061e1b; font-weight:750; margin-bottom:4px;'>Usuario:</p>", unsafe_allow_html=True)[cite: 16]
            u_in = st.text_input("Usuario", label_visibility="collapsed")[cite: 16]
            
            st.markdown("<p style='color:#061e1b; font-weight:750; margin-bottom:4px; margin-top:12px;'>PIN / Contraseña:</p>", unsafe_allow_html=True)[cite: 16]
            p_in = st.text_input("PIN", type="password", label_visibility="collapsed")[cite: 16]
            
            st.markdown("<div style='margin-top: 18px;'></div>", unsafe_allow_html=True)[cite: 16]
            if st.form_submit_button("Iniciar Sesión", use_container_width=True):
                u_clean = u_in.strip()
                if u_clean in db["usuarios"] and db["usuarios"][u_clean]["pin"] == p_in:
                    st.session_state.autenticado = True
                    st.session_state.usuario_clave = u_clean
                    st.rerun()
                else:
                    st.error("Credenciales incorrectas. Verifique el usuario y el PIN ingresado.")[cite: 16]
                    
        st.markdown("</div>", unsafe_allow_html=True)[cite: 16]
    st.stop()[cite: 16]

# --- DATOS DEL USUARIO ACTUAL ---
usr_key = st.session_state.usuario_clave
usr = db["usuarios"].get(usr_key, {"nombre": usr_key, "rol": "Usuario", "permiso_editar": False, "permiso_eliminar": False, "avatar": ""})
es_admin = (usr.get("rol") == "Admin") or (usr_key in ["Francesca Fellay", "gerente"])

# --- BARRA LATERAL ---
st.sidebar.markdown("### 👤 Mi Perfil")[cite: 16]

avatar_path = usr.get("avatar", "")[cite: 16]
if avatar_path and os.path.exists(avatar_path):
    st.sidebar.image(avatar_path, width=105)[cite: 16]
else:
    st.sidebar.markdown("""<div style='width:80px; height:80px; border-radius:50%; background:#a2d2cc; border:2px solid #52948d; display:flex; align-items:center; justify-content:center; font-size:2.2rem; color:#061e1b; margin-bottom:10px;'>👤</div>""", unsafe_allow_html=True)[cite: 16]

st.sidebar.markdown(f"**{usr.get('nombre')}**")[cite: 16]
st.sidebar.markdown(f"<span class='badge-role'>{usr.get('rol')}</span>", unsafe_allow_html=True)[cite: 16]

with st.sidebar.expander("📷 Cambiar foto de perfil"):
    nueva_foto = st.file_uploader("Subir imagen (JPG, PNG):", type=["jpg", "jpeg", "png"], key="upload_avatar")[cite: 16]
    if nueva_foto is not None:
        if st.button("Guardar Foto"):
            ruta_avatar = os.path.join(DIR_AVATARS, f"{usr_key}_avatar.png")[cite: 16]
            with open(ruta_avatar, "wb") as f_av:
                f_av.write(nueva_foto.getbuffer())[cite: 16]
            db["usuarios"][usr_key]["avatar"] = ruta_avatar[cite: 16]
            guardar_estado(db)[cite: 16]
            st.success("Foto actualizada.")[cite: 16]
            st.rerun()[cite: 16]

# --- GESTION CENTRAL DE COPIA DE SEGURIDAD (BACKUP UNIFICADO) ---
st.sidebar.markdown("---")[cite: 16]
st.sidebar.markdown("### 📦 Respaldo del Sistema")[cite: 16]

zip_backup_bytes = generar_zip_respaldo_completo()[cite: 16]
st.sidebar.download_button(
    label="📦 Descargar Backup Completo (.zip)",
    data=zip_backup_bytes,
    file_name=f"backup_plataforma_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip",
    mime="application/zip",
    use_container_width=True
)[cite: 16]

with st.sidebar.expander("🔄 Cargar / Restaurar Backup (.zip)"):
    archivo_restore = st.file_uploader("Subir archivo .zip de respaldo:", type=["zip"], key="up_zip_restore")[cite: 16]
    if archivo_restore is not None:
        if st.button("Restaurar Todo el Sistema", use_container_width=True):[cite: 16]
            restaurar_desde_zip_completo(archivo_restore.getvalue())[cite: 16]
            db = cargar_estado()[cite: 16]
            st.success("¡Base de datos y archivos restaurados con éxito!")[cite: 16]
            st.rerun()[cite: 16]

with st.sidebar.expander("📧 Enviar Backup al Correo"):[cite: 16]
    st.caption(f"Destino programado: **{CORREO_DESTINO_BACKUP}** (Automático a las 8:00 AM)")[cite: 16]
    if st.button("Enviar Respaldo Ahora"):[cite: 16]
        with st.spinner("Enviando respaldo por correo..."):[cite: 16]
            tz_cl = pytz.timezone("America/Santiago")[cite: 16]
            hora_fmt = datetime.now(tz_cl).strftime("%d/%m/%Y %H:%M:%S")[cite: 16]
            ok, msj = enviar_correo_backup(zip_backup_bytes, hora_fmt)[cite: 16]
            if ok:
                st.success("¡Correo de respaldo enviado exitosamente!")[cite: 16]
            else:
                st.error(f"Detalle: {msj}")[cite: 16]

st.sidebar.markdown("---")[cite: 16]
if st.sidebar.button("🚪 Cerrar Sesión", use_container_width=True):[cite: 16]
    st.session_state.autenticado = False[cite: 16]
    st.session_state.usuario_clave = None[cite: 16]
    st.rerun()[cite: 16]

# --- ENCABEZADO SUPERIOR CON RELOJ ---
col_head_title, col_head_clock = st.columns([1.2, 1])[cite: 16]
with col_head_title:
    st.markdown("<h1 style='color:#061e1b !important; margin:0;'>⚡ Centro de Gestión & Analítica</h1>", unsafe_allow_html=True)[cite: 16]
    st.markdown("<p style='color:#061e1b; margin-top:4px;'>Plataforma colaborativa multi-formato con procesamiento inteligente mediante Groq IA.</p>", unsafe_allow_html=True)[cite: 16]
with col_head_clock:
    renderizar_reloj_chile()[cite: 16]

# --- NAVEGACION POR PESTAÑAS ---
titulos_pestanas = ["📊 Datasets & Archivos", "📂 Material", "📄 Informes Compartidos"][cite: 16]
if es_admin:
    titulos_pestanas.append("👥 Gestión de Usuarios")[cite: 16]

pestanas_principales = st.tabs(titulos_pestanas)[cite: 16]

# ==========================================
# 1. DATASETS EXCEL
# ==========================================
with pestanas_principales[0]:[cite: 16]
    st.markdown("""<div class='modern-card'><h3 style='margin:0;'>📊 Repositorio de Datasets</h3><p style='margin:0; color:#061e1b;'>Suba, visualice y edite hojas de cálculo en ventanas independientes.</p></div>""", unsafe_allow_html=True)[cite: 16]
    
    with st.expander("➕ Cargar Nuevo Dataset Excel", expanded=True):[cite: 16]
        c_tit, c_arc = st.columns([1, 1])[cite: 16]
        with c_tit:
            t_data = st.text_input("Título del dataset:")[cite: 16]
        with c_arc:
            f_data = st.file_uploader("Archivo Excel (.xlsx, .xls):", type=["xlsx", "xls"])[cite: 16]
            
        if st.button("Guardar Dataset"):[cite: 16]
            if not t_data.strip() or f_data is None:
                st.error("Complete el título y seleccione un archivo.")[cite: 16]
            else:
                nom_arc = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{f_data.name}"[cite: 16]
                ruta_dest = os.path.join(DIR_DATASETS, nom_arc)[cite: 16]
                with open(ruta_dest, "wb") as f:
                    f.write(f_data.getbuffer())[cite: 16]
                
                db["datasets"][t_data] = {
                    "titulo": t_data,
                    "autor": usr["nombre"],
                    "fecha": datetime.now().strftime("%Y-%m-%d %H:%M"),
                    "ruta": ruta_dest
                }[cite: 16]
                guardar_estado(db)[cite: 16]
                st.success("Dataset guardado con éxito.")[cite: 16]
                st.rerun()[cite: 16]

    st.markdown("---")[cite: 16]
    if not db["datasets"]:
        st.info("No hay datasets subidos actualmente.")[cite: 16]
    else:
        titulos_ds = list(db["datasets"].keys())[cite: 16]
        pestanas_ds = st.tabs(titulos_ds)[cite: 16]
        
        for idx_ds, tab_ds in enumerate(pestanas_ds):[cite: 16]
            t_act = titulos_ds[idx_ds][cite: 16]
            info_ds = db["datasets"][t_act][cite: 16]
            
            with tab_ds:
                st.markdown(f"### 📋 {info_ds['titulo']}")[cite: 16]
                st.caption(f"Subido por: **{info_ds['autor']}** | Fecha: {info_ds['fecha']}")[cite: 16]
                
                if os.path.exists(info_ds["ruta"]):
                    df_actual = pd.read_excel(info_ds["ruta"])[cite: 16]
                    
                    if usr.get("permiso_editar") or es_admin:[cite: 16]
                        st.markdown("**✏️ Editor de Datos en Vivo:**")[cite: 16]
                        df_edit = st.data_editor(df_actual, key=f"d_edit_{t_act}", use_container_width=True)[cite: 16]
                        if st.button("💾 Guardar Cambios en Excel", key=f"s_df_{t_act}"):[cite: 16]
                            df_edit.to_excel(info_ds["ruta"], index=False)[cite: 16]
                            st.success("Datos actualizados.")[cite: 16]
                            st.rerun()[cite: 16]
                    else:
                        st.dataframe(df_actual, use_container_width=True)[cite: 16]
                    
                    if usr.get("permiso_eliminar") or es_admin:[cite: 16]
                        st.markdown("---")[cite: 16]
                        if st.button("🗑️ Eliminar Dataset", key=f"del_ds_{t_act}", type="secondary"):[cite: 16]
                            if os.path.exists(info_ds["ruta"]):
                                os.remove(info_ds["ruta"])[cite: 16]
                            del db["datasets"][t_act][cite: 16]
                            guardar_estado(db)[cite: 16]
                            st.success("Dataset eliminado.")[cite: 16]
                            st.rerun()[cite: 16]
                else:
                    st.error("Archivo físico no encontrado.")[cite: 16]

# ==========================================
# 2. MODULO MATERIAL MULTI-FORMATO
# ==========================================
with pestanas_principales[1]:[cite: 16]
    st.markdown("""<div class='modern-card'><h3 style='margin:0;'>📂 Módulo de Material</h3><p style='margin:0; color:#061e1b;'>Soporte y resúmenes automáticos con Groq para Documentos, Imágenes y Audios (M4A/MP3/WAV).</p></div>""", unsafe_allow_html=True)[cite: 16]
    
    with st.form("form_material"):[cite: 16]
        tit_mat = st.text_input("Título descriptivo del material:")[cite: 16]
        archivos = st.file_uploader(
            "Cargar archivos:",
            type=["xlsx", "xls", "docx", "pdf", "pptx", "jpg", "jpeg", "png", "mp4", "m4a", "mp3", "wav"],
            accept_multiple_files=True
        )[cite: 16]
        if st.form_submit_button("Subir y Procesar"):[cite: 16]
            if not tit_mat.strip() or not archivos:
                st.error("Complete el título y cargue al menos un archivo.")[cite: 16]
            else:
                client = obtener_cliente_ia()[cite: 16]
                for a in archivos:
                    ext = a.name.split(".")[-1].lower()[cite: 16]
                    nom_dest = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{a.name}"[cite: 16]
                    ruta_guardada = os.path.join(DIR_MATERIAL, nom_dest)[cite: 16]
                    
                    with open(ruta_guardada, "wb") as f:
                        f.write(a.getbuffer())[cite: 16]
                    
                    resumen_txt = "Archivo guardado."[cite: 16]
                    if client:
                        with st.spinner(f"Analizando {a.name} con Groq IA..."):[cite: 16]
                            try:
                                if ext in ["m4a", "mp3", "wav"]:
                                    with open(ruta_guardada, "rb") as f_aud:
                                        audio_bytes = f_aud.read()[cite: 16]
                                    audio_buffer = io.BytesIO(audio_bytes)[cite: 16]
                                    audio_buffer.name = f"temp_audio.{ext}"[cite: 16]
                                    
                                    transcripcion = client.audio.transcriptions.create(
                                        model=MODELO_WHISPER,
                                        file=audio_buffer,
                                        language="es"
                                    )[cite: 16]
                                    texto_audio = transcripcion.text[cite: 16]
                                    
                                    p_sys = "Eres un especialista senior en diagnóstico y análisis de materiales. Sintetiza con precisión en español latinoamericano."[cite: 16]
                                    p_user = f"A partir de la siguiente transcripción de audio ({a.name}), sintetiza los puntos clave tratados, acuerdos y diagnóstico. Usa párrafos claros o viñetas Markdown estándar, nunca diagramas ASCII:\n\n{texto_audio}"[cite: 16]
                                    resumen_txt = ejecutar_chat_groq(client, p_sys, p_user)[cite: 16]
                                elif ext in ["pdf", "docx", "pptx", "xlsx", "xls"]:
                                    t_doc = extraer_texto_archivo(ruta_guardada, ext)[cite: 16]
                                    p_sys = "Eres un especialista senior en análisis documental y materiales."[cite: 16]
                                    p_user = f"Elabora un resumen y diagnóstico clave del siguiente documento ({a.name}):\n\n{t_doc}\n\nUsa viñetas Markdown estándar con conceptos clave en negrita."[cite: 16]
                                    resumen_txt = ejecutar_chat_groq(client, p_sys, p_user)[cite: 16]
                                elif ext in ["jpg", "jpeg", "png"]:
                                    resumen_txt = f"Imagen registrada ({a.name}). Visualización disponible en la pestaña multimedia."[cite: 16]
                                elif ext == "mp4":
                                    resumen_txt = f"Video registrado ({a.name}). Reproducción multimedia disponible."[cite: 16]
                            except Exception as e:
                                resumen_txt = f"Archivo guardado. Diagnóstico no generado: {e}"[cite: 16]
                    
                    db["material"].append({
                        "titulo": tit_mat,
                        "nombre_original": a.name,
                        "tipo": ext,
                        "autor": usr["nombre"],
                        "fecha": datetime.now().strftime("%Y-%m-%d %H:%M"),
                        "ruta": ruta_guardada,
                        "resumen": resumen_txt
                    })[cite: 16]
                guardar_estado(db)[cite: 16]
                st.success("Materiales integrados correctamente.")[cite: 16]
                st.rerun()[cite: 16]

    st.markdown("---")[cite: 16]
    st.subheader("📚 Materiales Guardados")[cite: 16]
    
    if not db["material"]:
        st.info("No hay registros en material.")[cite: 16]
    else:
        for idx, item in enumerate(db["material"]):[cite: 16]
            titulo_item = item.get("titulo") or item.get("titulo_registro") or "Sin Título"[cite: 16]
            nombre_arc = item.get("nombre_original") or item.get("filename") or "Archivo"[cite: 16]
            tipo_arc = item.get("tipo", "").lower()[cite: 16]
            autor_arc = item.get("autor", "Desconocido")[cite: 16]
            fecha_arc = item.get("fecha", "")[cite: 16]
            resumen_arc = item.get("resumen", "Sin resumen disponible.")[cite: 16]
            ruta_arc = item.get("ruta", "")[cite: 16]
            
            with st.container():[cite: 16]
                col_head, col_del_btn = st.columns([5, 1])[cite: 16]
                with col_head:
                    st.markdown(f"### 📁 {titulo_item} — <span class='highlight-tag'>{nombre_arc}</span>", unsafe_allow_html=True)[cite: 16]
                    st.caption(f"Autor: **{autor_arc}** | Fecha: {fecha_arc} | Formato: **{tipo_arc.upper()}**")[cite: 16]
                with col_del_btn:
                    if usr.get("permiso_eliminar") or es_admin:[cite: 16]
                        if st.button("🗑️ Borrar", key=f"del_mat_{idx}", type="secondary"):[cite: 16]
                            if ruta_arc and os.path.exists(ruta_arc):
                                os.remove(ruta_arc)[cite: 16]
                            if os.path.exists(os.path.join(DIR_USUARIOS, f"{nombre_arc}.json")):
                                os.remove(os.path.join(DIR_USUARIOS, f"{nombre_arc}.json"))[cite: 16]
                            db["material"].pop(idx)[cite: 16]
                            guardar_estado(db)[cite: 16]
                            st.success("Archivo eliminado.")[cite: 16]
                            st.rerun()[cite: 16]
                
                t_vis, t_res = st.tabs(["👁️ Multimedia / Descarga", "📝 Diagnóstico y Resumen"])[cite: 16]
                with t_vis:
                    if ruta_arc and os.path.exists(ruta_arc):
                        if tipo_arc == "mp4":
                            st.video(ruta_arc)[cite: 16]
                        elif tipo_arc in ["m4a", "mp3", "wav"]:
                            st.audio(ruta_arc)[cite: 16]
                        elif tipo_arc in ["jpg", "jpeg", "png"]:
                            st.image(ruta_arc, use_container_width=True)[cite: 16]
                        else:
                            with open(ruta_arc, "rb") as fl:
                                st.download_button("📥 Descargar Archivo", data=fl.read(), file_name=nombre_arc, key=f"dl_a_{idx}")[cite: 16]
                    else:
                        st.error("Archivo físico no encontrado.")[cite: 16]
                with t_res:
                    st.markdown(resumen_arc)[cite: 16]
                st.markdown("---")[cite: 16]

# ==========================================
# 3. INFORMES COMPARTIDOS
# ==========================================
with pestanas_principales[2]:[cite: 16]
    st.markdown("""<div class='modern-card'><h3 style='margin:0;'>📄 Informes Compartidos</h3><p style='margin:0; color:#061e1b;'>Generación y repositorio colaborativo con exportación en formato Carta.</p></div>""", unsafe_allow_html=True)[cite: 16]
    
    with st.expander("🤖 Redactar Nuevo Informe con IA (Groq)", expanded=False):[cite: 16]
        nom_i = st.text_input("Título del Informe:")[cite: 16]
        enf_i = st.selectbox("Enfoque:", ["Resumen Ejecutivo", "Diagnóstico Técnico", "Evaluación Estratégica"])[cite: 16]
        ins_i = st.text_area("Instrucciones complementarias:")[cite: 16]
        
        if st.button("🚀 Generar Informe"):[cite: 16]
            client = obtener_cliente_ia()[cite: 16]
            if not nom_i.strip():
                st.error("Ingrese un título para el informe.")[cite: 16]
            elif not client:
                st.error("API Key de Groq (GROQ_API_KEY) no configurada en Secrets.")[cite: 16]
            else:
                with st.spinner("Redactando informe consolidado con Groq IA..."):[cite: 16]
                    resumenes_mat = "\n".join([f"- {x.get('nombre_original', 'Archivo')}: {x.get('resumen', '')}" for x in db["material"]])[cite: 16]
                    
                    p_sys = "Actúa como especialista analítico senior. Genera informes estructurados de excelencia en español latinoamericano."[cite: 16]
                    p_user = f"""Genera un informe estructurado profesional con enfoque '{enf_i}'.
Título: {nom_i}
Autor solicitante: {usr['nombre']}
Instrucciones: {ins_i}
Materiales analizados: {resumenes_mat if resumenes_mat else 'Sin materiales adjuntos.'}

ESTRUCTURA REQUERIDA:
1. Diagnóstico y Visión Global
2. Hallazgos Analíticos Relevantes (Si incluyes tablas comparativas o síntesis, hazlo OBLIGATORIAMENTE usando sintaxis estándar de tablas Markdown con '|' y encabezados separados por '|---|---|', NUNCA uses arte ASCII ni bloques de código con ``` para tablas).
3. Conclusiones y Plan de Acción

REGLA OBLIGATORIA DE FORMATO:
- Las tablas deben ser tablas Markdown renderizables reales.
- Todo el texto analítico debe fluir con negritas y viñetas claras."""[cite: 16]
                    try:
                        resp_txt = ejecutar_chat_groq(client, p_sys, p_user)[cite: 16]
                        
                        nuevo_inf = {
                            "titulo": nom_i,
                            "autor": usr["nombre"],
                            "fecha": datetime.now().strftime("%Y-%m-%d %H:%M"),
                            "contenido": resp_txt
                        }[cite: 16]
                        
                        db["informes"].append(nuevo_inf)[cite: 16]
                        
                        nom_inf_f = f"inf_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"[cite: 16]
                        with open(os.path.join(DIR_INFORMES, nom_inf_f), "w", encoding="utf-8") as fi_b:
                            json.dump(nuevo_inf, fi_b, ensure_ascii=False, indent=2)[cite: 16]
                            
                        guardar_estado(db)[cite: 16]
                        st.success("Informe publicado y respaldado con éxito.")[cite: 16]
                        st.rerun()[cite: 16]
                    except Exception as err:
                        st.error(f"Error al generar informe: {err}")[cite: 16]

    st.markdown("---")[cite: 16]
    st.subheader("📚 Repositorio de Informes")[cite: 16]
    
    if not db["informes"]:
        st.info("No hay informes registrados.")[cite: 16]
    else:
        for idx_inf, inf in enumerate(db["informes"]):[cite: 16]
            titulo_inf = inf.get("titulo") or inf.get("titulo_informe") or "Sin Título"[cite: 16]
            autor_inf = inf.get("autor", "Desconocido")[cite: 16]
            fecha_inf = inf.get("fecha", "")[cite: 16]
            contenido_inf = inf.get("contenido", "")[cite: 16]
            
            with st.container():[cite: 16]
                c_inf_t, c_inf_del = st.columns([5, 1])[cite: 16]
                with c_inf_t:
                    st.markdown(f"### 📄 [{autor_inf}] - {titulo_inf}")[cite: 16]
                    st.caption(f"Generado el: {fecha_inf}")[cite: 16]
                with c_inf_del:
                    if usr.get("permiso_eliminar") or es_admin:[cite: 16]
                        if st.button("🗑️ Borrar", key=f"del_inf_{idx_inf}", type="secondary"):[cite: 16]
                            db["informes"].pop(idx_inf)[cite: 16]
                            guardar_estado(db)[cite: 16]
                            st.success("Informe eliminado.")[cite: 16]
                            st.rerun()[cite: 16]
                
                st.markdown(contenido_inf)[cite: 16]
                
                cd1, cd2 = st.columns([1, 1])[cite: 16]
                with cd1:
                    df_out = pd.DataFrame([{"Autor": autor_inf, "Título": titulo_inf, "Fecha": fecha_inf, "Contenido": contenido_inf}])[cite: 16]
                    buf = io.BytesIO()[cite: 16]
                    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
                        df_out.to_excel(writer, index=False)[cite: 16]
                    st.download_button(
                        label="📥 Descargar Excel (.xlsx)",
                        data=buf.getvalue(),
                        file_name=f"[{autor_inf}] - {titulo_inf}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        key=f"xls_dl_{idx_inf}"
                    )[cite: 16]
                with cd2:
                    html_carta = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{titulo_inf}</title>
    <style>
        @page {{ size: letter portrait; margin: 25mm; }}
        body {{ font-family: 'Segoe UI', Arial, sans-serif; color: #061e1b; padding: 20px; background-color: #89c7c0; }}
        .header {{ border-bottom: 2px solid #52948d; padding-bottom: 8px; margin-bottom: 20px; }}
        .title {{ font-size: 20pt; font-weight: bold; color: #061e1b; }}
        .meta {{ font-size: 10pt; color: #061e1b; margin-top: 4px; }}
        .body {{ font-size: 11pt; line-height: 1.6; white-space: pre-wrap; }}
        table {{ width: 100%; border-collapse: collapse; margin: 15px 0; }}
        th, td {{ border: 1px solid #52948d; padding: 8px 12px; text-align: left; }}
        th {{ background-color: #63a59e; }}
    </style>
</head>
<body>
    <div class="header">
        <div class="title">{titulo_inf}</div>
        <div class="meta">Autor: {autor_inf} | Fecha: {fecha_inf}</div>
    </div>
    <div class="body">{contenido_inf}</div>
</body>
</html>"""[cite: 16]
                    st.download_button(
                        label="🖼️ Descargar Formato Carta (HTML / PDF)",
                        data=html_carta,
                        file_name=f"[{autor_inf}] - {titulo_inf}_Carta.html",
                        mime="text/html",
                        key=f"doc_dl_{idx_inf}"
                    )[cite: 16]
                st.markdown("---")[cite: 16]

# ==========================================
# 4. GESTION DE USUARIOS (SOLO ADMIN)
# ==========================================
if es_admin:
    with pestanas_principales[3]:[cite: 16]
        st.markdown("""<div class='modern-card'><h3 style='margin:0;'>👥 Gestión de Usuarios y Permisos</h3><p style='margin:0; color:#061e1b;'>Control centralizado de cuentas accesible exclusivamente por Administradores.</p></div>""", unsafe_allow_html=True)[cite: 16]
        
        with st.expander("➕ Registrar Nuevo Usuario", expanded=True):[cite: 16]
            cu1, cu2, cu3 = st.columns([1, 1, 1])[cite: 16]
            with cu1:
                n_u = st.text_input("Usuario (Login):")[cite: 16]
                n_nom = st.text_input("Nombre Completo:")[cite: 16]
            with cu2:
                n_pin = st.text_input("PIN / Clave:", type="password")[cite: 16]
                n_rol = st.selectbox("Rol:", ["Usuario", "Analista", "Especialista", "Gerencia", "Admin"])[cite: 16]
            with cu3:
                st.markdown("**Permisos Iniciales:**")[cite: 16]
                p_e = st.checkbox("Permiso para Editar Datasets")[cite: 16]
                p_d = st.checkbox("Permiso para Eliminar Datasets / Archivos")[cite: 16]
                
            if st.button("Crear Usuario"):[cite: 16]
                if not n_u or not n_pin or not n_nom:
                    st.error("Todos los campos son obligatorios.")[cite: 16]
                elif n_u in db["usuarios"] or n_u == "gerente":
                    st.error("El nombre de usuario ya existe.")[cite: 16]
                else:
                    db["usuarios"][n_u] = {
                        "nombre": n_nom,
                        "rol": n_rol,
                        "pin": n_pin,
                        "permiso_editar": p_e,
                        "permiso_eliminar": p_d,
                        "avatar": ""
                    }[cite: 16]
                    guardar_estado(db)[cite: 16]
                    st.success(f"Usuario '{n_u}' registrado y respaldado correctamente.")[cite: 16]
                    st.rerun()[cite: 16]

        st.markdown("---")[cite: 16]
        st.subheader("📜 Cuentas Registradas")[cite: 16]
        
        # Filtro estricto para mantener al usuario gerente completamente oculto
        usuarios_visibles = [k for k in list(db["usuarios"].keys()) if k != "gerente"][cite: 16]
        
        for u_k in usuarios_visibles:[cite: 16]
            u_d = db["usuarios"][u_k][cite: 16]
            with st.container():[cite: 16]
                col_u_inf, col_u_e, col_u_d, col_u_del = st.columns([2, 1, 1, 1])[cite: 16]
                with col_u_inf:
                    st.markdown(f"**{u_d.get('nombre', '')}** (<span class='highlight-tag'>{u_k}</span>) — Rol: *{u_d.get('rol', '')}*", unsafe_allow_html=True)[cite: 16]
                
                if u_k == "Francesca Fellay":
                    st.caption("Administrador Principal (Cuenta protegida)")[cite: 16]
                else:
                    with col_u_e:
                        val_e = st.checkbox("Editar", value=u_d.get("permiso_editar", False), key=f"pe_{u_k}")[cite: 16]
                    with col_u_d:
                        val_d = st.checkbox("Eliminar", value=u_d.get("permiso_eliminar", False), key=f"pd_{u_k}")[cite: 16]
                    with col_u_del:
                        if st.button("🗑️ Eliminar", key=f"del_user_{u_k}", type="secondary"):[cite: 16]
                            ruta_u_del = os.path.join(DIR_USUARIOS, f"{u_k}.json")[cite: 16]
                            if os.path.exists(ruta_u_del):
                                os.remove(ruta_u_del)[cite: 16]
                            del db["usuarios"][u_k][cite: 16]
                            guardar_estado(db)[cite: 16]
                            st.success(f"Usuario {u_k} eliminado.")[cite: 16]
                            st.rerun()[cite: 16]
                            
                    if val_e != u_d.get("permiso_editar") or val_d != u_d.get("permiso_eliminar"):[cite: 16]
                        db["usuarios"][u_k]["permiso_editar"] = val_e[cite: 16]
                        db["usuarios"][u_k]["permiso_eliminar"] = val_d[cite: 16]
                        guardar_estado(db)[cite: 16]
            st.markdown("---")[cite: 16]
