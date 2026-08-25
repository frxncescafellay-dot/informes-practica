import os
import io
import json
from datetime import datetime
import pytz
import pandas as pd
import streamlit as st
import streamlit.components.v1 as components
from PIL import Image
from groq import Groq
import pypdf
from docx import Document
from pptx import Presentation

# --- CONFIGURACION DE PAGINA ---
st.set_page_config(
    page_title="Plataforma de Analitica & Intervencion IA",
    page_icon="⚡",
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- RUTAS DE ALMACENAMIENTO PERSISTENTE ---
DIR_BASE = "almacen_datos"
DIR_DATASETS = os.path.join(DIR_BASE, "datasets")
DIR_INTERVENCION = os.path.join(DIR_BASE, "intervencion")
DIR_AVATARS = os.path.join(DIR_BASE, "avatares")
FILE_DB = os.path.join(DIR_BASE, "base_datos.json")

for d in [DIR_BASE, DIR_DATASETS, DIR_INTERVENCION, DIR_AVATARS]:
    os.makedirs(d, exist_ok=True)[cite: 7]

MODELO_WHISPER = "whisper-large-v3"[cite: 3, 6]

# --- PALETA TEAL + NARANJA TERRACOTA (CERO ELEMENTOS NEGROS) ---
st.markdown("""
<style>
    /* Tipografia uniforme */
    html, body, [class*="css"], .stApp {
        font-family: 'Segoe UI', system-ui, -apple-system, sans-serif !important;
    }

    /* Fondo principal */
    .stApp {
        background: #73b5ae !important;
        color: #061e1b;
    }
    
    /* Encabezados y titulos */
    h1, h2, h3, h4, h5, h6 {
        color: #061e1b !important;
        font-weight: 800;
        letter-spacing: -0.5px;
        font-family: 'Segoe UI', system-ui, -apple-system, sans-serif !important;
    }
    
    /* Barra lateral */
    section[data-testid="stSidebar"] {
        background-color: #63a59e !important;
        border-right: 2px solid #52948d !important;
    }
    section[data-testid="stSidebar"] * {
        color: #061e1b !important;
        font-weight: 600;
    }

    /* Tarjetas y Contenedores */
    .modern-card {
        background: #89c7c0 !important;
        border: 2px solid #63a59e !important;
        border-radius: 14px;
        padding: 22px;
        margin-bottom: 20px;
        box-shadow: 0 4px 12px rgba(13, 44, 41, 0.12);
    }

    /* Desplegables (Expanders) */
    div[data-testid="stExpander"] {
        background-color: #89c7c0 !important;
        border: 2px solid #63a59e !important;
        border-radius: 10px !important;
    }
    div[data-testid="stExpander"] summary {
        background-color: #7ab8b1 !important;
        color: #061e1b !important;
        border-radius: 8px !important;
        font-weight: 750 !important;
    }
    div[data-testid="stExpander"] summary * {
        color: #061e1b !important;
        font-weight: 750 !important;
    }
    div[data-testid="stExpander"] div[role="region"] {
        background-color: #89c7c0 !important;
        color: #061e1b !important;
    }

    /* Cuadros de texto */
    input[type="text"], 
    input[type="password"], 
    textarea {
        background-color: #a2d2cc !important;
        color: #061e1b !important;
        border: 2px solid #52948d !important;
        border-radius: 8px !important;
        font-size: 0.95rem !important;
        font-weight: 600 !important;
    }
    input[type="text"]:focus, 
    input[type="password"]:focus, 
    textarea:focus {
        background-color: #bfe3de !important;
        color: #000000 !important;
        border-color: #c84b1e !important;
        box-shadow: 0 0 0 3px rgba(200, 75, 30, 0.25) !important;
    }

    /* Quitar fondo negro del boton de ver contrasena en input password */
    div[data-baseweb="input"] button,
    div[data-baseweb="input"] div {
        background-color: #a2d2cc !important;
        border: none !important;
        color: #061e1b !important;
    }
    div[data-baseweb="input"] svg {
        fill: #061e1b !important;
        color: #061e1b !important;
    }

    /* Selectores desplegables (Selectbox) */
    div[data-baseweb="select"] > div {
        background-color: #a2d2cc !important;
        color: #061e1b !important;
        border: 2px solid #52948d !important;
        border-radius: 8px !important;
    }
    div[data-baseweb="select"] * {
        background-color: #a2d2cc !important;
        color: #061e1b !important;
        font-weight: 600 !important;
    }
    div[data-baseweb="select"] svg {
        fill: #061e1b !important;
        color: #061e1b !important;
    }

    /* CUADROS PARA MARCAR (Checkboxes) */
    div[data-testid="stCheckbox"] span[role="checkbox"] {
        background-color: #a2d2cc !important;
        border: 2px solid #52948d !important;
        border-radius: 5px !important;
    }
    div[data-testid="stCheckbox"] span[role="checkbox"][aria-checked="true"] {
        background-color: #c84b1e !important;
        border-color: #9e3610 !important;
    }
    div[data-testid="stCheckbox"] svg {
        color: #ffffff !important;
        fill: #ffffff !important;
    }
    div[data-testid="stCheckbox"] label span {
        color: #061e1b !important;
        font-weight: 600 !important;
    }

    /* Zona de carga de archivos (Uploader) */
    div[data-testid="stFileUploader"] {
        background-color: #89c7c0 !important;
        border: 2px dashed #c84b1e !important;
        border-radius: 12px !important;
        padding: 14px !important;
    }
    div[data-testid="stFileUploader"] section {
        background-color: #89c7c0 !important;
    }
    div[data-testid="stFileUploader"] section * {
        color: #061e1b !important;
    }
    div[data-testid="stFileUploader"] button {
        background-color: #c84b1e !important;
        border: 2px solid #9e3610 !important;
        border-radius: 8px !important;
        font-weight: 750 !important;
    }
    div[data-testid="stFileUploader"] button p {
        color: #ffffff !important;
    }

    /* BOTONES DE DESCARGA */
    div[data-testid="stDownloadButton"]>button {
        background: #a2d2cc !important;
        color: #061e1b !important;
        border: 2px solid #52948d !important;
        border-radius: 9px !important;
        padding: 9px 20px !important;
        font-weight: 800 !important;
        font-size: 0.95rem !important;
        box-shadow: 0 3px 8px rgba(13, 44, 41, 0.15) !important;
        transition: all 0.2s ease !important;
    }
    div[data-testid="stDownloadButton"]>button:hover {
        background: #bfe3de !important;
        color: #000000 !important;
        border-color: #061e1b !important;
        transform: translateY(-1px);
    }
    div[data-testid="stDownloadButton"]>button p {
        color: #061e1b !important;
        font-weight: 800 !important;
    }

    /* BOTONES PRINCIPALES Y FORMULARIOS */
    .stButton>button, 
    div[data-testid="stFormSubmitButton"]>button {
        background: linear-gradient(135deg, #c84b1e 0%, #b23b14 100%) !important;
        color: #ffffff !important;
        border: 2px solid #9e3610 !important;
        border-radius: 9px !important;
        padding: 9px 22px !important;
        font-weight: 800 !important;
        font-size: 0.95rem !important;
        box-shadow: 0 4px 12px rgba(178, 59, 20, 0.35) !important;
        transition: all 0.2s ease !important;
    }
    .stButton>button:hover, 
    div[data-testid="stFormSubmitButton"]>button:hover {
        background: linear-gradient(135deg, #db5928 0%, #c84b1e 100%) !important;
        border-color: #db5928 !important;
        box-shadow: 0 6px 16px rgba(219, 89, 40, 0.45) !important;
        transform: translateY(-1px);
    }
    .stButton>button p,
    div[data-testid="stFormSubmitButton"]>button p {
        color: #ffffff !important;
        font-weight: 800 !important;
    }

    /* Tablas de datos Excel */
    div[data-testid="stDataFrame"], 
    div[data-testid="stDataEditor"],
    div[data-testid="stDataFrame"] > div,
    div[data-testid="stDataEditor"] > div {
        background-color: #89c7c0 !important;
        border: 2px solid #52948d !important;
        border-radius: 10px !important;
    }
    
    .glide-data-grid,
    .gdg-container,
    div[data-testid="stDataFrame"] canvas,
    div[data-testid="stDataEditor"] canvas {
        filter: invert(0.85) sepia(0.3) saturate(2.5) hue-rotate(130deg) brightness(1.05) !important;
        border-radius: 8px !important;
    }

    /* TABLAS MARKDOWN EN INFORMES */
    table {
        width: 100% !important;
        border-collapse: collapse !important;
        margin: 16px 0 !important;
        background-color: #89c7c0 !important;
        border-radius: 8px !important;
        overflow: hidden !important;
        border: 2px solid #52948d !important;
    }
    th {
        background-color: #63a59e !important;
        color: #061e1b !important;
        font-weight: 800 !important;
        padding: 12px 14px !important;
        border: 1px solid #52948d !important;
        text-align: left !important;
    }
    td {
        padding: 10px 14px !important;
        border: 1px solid #52948d !important;
        color: #061e1b !important;
        font-size: 0.95rem !important;
    }
    tr:nth-child(even) {
        background-color: #7ab8b1 !important;
    }

    /* Textos con fondo destacado */
    .highlight-tag {
        background: #a2d2cc;
        border: 1.5px solid #52948d;
        padding: 3px 10px;
        border-radius: 6px;
        font-weight: 700;
        color: #061e1b;
        font-size: 0.92rem;
        display: inline-block;
        font-family: 'Segoe UI', system-ui, -apple-system, sans-serif !important;
    }

    code, pre {
        background-color: #a2d2cc !important;
        color: #061e1b !important;
        border: 1.5px solid #52948d !important;
        border-radius: 6px !important;
        padding: 3px 8px !important;
        font-family: 'Segoe UI', system-ui, -apple-system, sans-serif !important;
        font-size: 0.92rem !important;
        font-weight: 700 !important;
    }
    
    /* Botones de eliminacion */
    div[data-testid="stBaseButton-secondary"] button {
        background: #e29b9b !important;
        border: 2px solid #c96b6b !important;
        box-shadow: none !important;
    }
    div[data-testid="stBaseButton-secondary"] button:hover {
        background: #d67a7a !important;
    }
    div[data-testid="stBaseButton-secondary"] button p {
        color: #5c0f0f !important;
    }

    /* Pestanas de navegacion */
    button[data-baseweb="tab"] {
        background-color: #63a59e !important;
        color: #061e1b !important;
        border-radius: 8px 8px 0 0 !important;
        font-weight: 750 !important;
        margin-right: 4px !important;
        padding: 9px 18px !important;
        border: 1px solid #52948d !important;
    }
    button[data-baseweb="tab"][aria-selected="true"] {
        background-color: #89c7c0 !important;
        color: #061e1b !important;
        border-bottom: 3.5px solid #c84b1e !important;
    }

    label, p, span {
        color: #061e1b;
        font-weight: 600;
    }

    .badge-role {
        display: inline-block;
        background: #c84b1e;
        color: #ffffff !important;
        padding: 4px 12px;
        border-radius: 14px;
        font-size: 0.82rem;
        font-weight: 800;
        border: 1.5px solid #9e3610;
    }
</style>
""", unsafe_allow_html=True)[cite: 7]

# --- RELOJ DUAL NARANJA ---
def renderizar_reloj_chile():
    html_reloj = """
    <div style="background: linear-gradient(135deg, #c84b1e 0%, #b23b14 100%); border: 2px solid #9e3610; border-radius: 16px; padding: 12px 20px; display: flex; align-items: center; justify-content: space-between; gap: 16px; box-shadow: 0 6px 18px rgba(178, 59, 20, 0.35); max-width: 480px; margin-left: auto; margin-bottom: 15px; font-family: 'Segoe UI', system-ui, sans-serif;">
        <div style="position: relative; width: 60px; height: 60px;">
            <svg id="analog-clock" width="60" height="60" viewBox="0 0 100 100">
                <circle cx="50" cy="50" r="46" fill="#fbe8e1" stroke="#782307" stroke-width="4"/>
                <line x1="50" y1="10" x2="50" y2="16" stroke="#4a1503" stroke-width="3" stroke-linecap="round"/>
                <line x1="90" y1="50" x2="84" y2="50" stroke="#4a1503" stroke-width="3" stroke-linecap="round"/>
                <line x1="50" y1="90" x2="50" y2="84" stroke="#4a1503" stroke-width="3" stroke-linecap="round"/>
                <line x1="10" y1="50" x2="16" y2="50" stroke="#4a1503" stroke-width="3" stroke-linecap="round"/>
                <line id="hour-hand" x1="50" y1="50" x2="50" y2="28" stroke="#260900" stroke-width="5" stroke-linecap="round"/>
                <line id="min-hand" x1="50" y1="50" x2="50" y2="18" stroke="#591905" stroke-width="3.5" stroke-linecap="round"/>
                <line id="sec-hand" x1="50" y1="50" x2="50" y2="14" stroke="#b23b14" stroke-width="2" stroke-linecap="round"/>
                <circle cx="50" cy="50" r="4" fill="#260900"/>
            </svg>
        </div>
        <div style="display: flex; flex-direction: column; justify-content: center;">
            <div style="font-size: 0.76rem; font-weight: 800; color: #fbe8e1; text-transform: uppercase; letter-spacing: 0.8px;">🇨🇱 Hora Oficial de Chile</div>
            <div id="digital-clock" style="font-size: 1.65rem; font-weight: 900; color: #ffffff; line-height: 1.1; text-shadow: 0 2px 4px rgba(0,0,0,0.2);">--:--:--</div>
            <div id="digital-date" style="font-size: 0.82rem; font-weight: 700; color: #fbe8e1; margin-top: 2px;">Cargando fecha...</div>
        </div>
    </div>
    <script>
        function actualizarReloj() {
            const opciones = { timeZone: 'America/Santiago', hour12: false, hour: '2-digit', minute: '2-digit', second: '2-digit' };
            const opcionesFecha = { timeZone: 'America/Santiago', weekday: 'long', year: 'numeric', month: 'long', day: 'numeric' };
            const ahora = new Date();
            const formateadorHora = new Intl.DateTimeFormat('es-CL', opciones);
            const formateadorFecha = new Intl.DateTimeFormat('es-CL', opcionesFecha);
            const partesHora = formateadorHora.formatToParts(ahora);
            let h = 0, m = 0, s = 0;
            partesHora.forEach(p => {
                if (p.type === 'hour') h = parseInt(p.value);
                if (p.type === 'minute') m = parseInt(p.value);
                if (p.type === 'second') s = parseInt(p.value);
            });
            const hStr = h.toString().padStart(2, '0');
            const mStr = m.toString().padStart(2, '0');
            const sStr = s.toString().padStart(2, '0');
            document.getElementById('digital-clock').textContent = `${hStr}:${mStr}:${sStr}`;
            let fechaStr = formateadorFecha.format(ahora);
            fechaStr = fechaStr.charAt(0).toUpperCase() + fechaStr.slice(1);
            document.getElementById('digital-date').textContent = fechaStr;
            const secDeg = s * 6;
            const minDeg = m * 6 + s * 0.1;
            const hourDeg = (h % 12) * 30 + m * 0.5;
            function rotar(elemId, deg) {
                const el = document.getElementById(elemId);
                if (el) el.setAttribute('transform', `rotate(${deg} 50 50)`);
            }
            rotar('sec-hand', secDeg);
            rotar('min-hand', minDeg);
            rotar('hour-hand', hourDeg);
        }
        setInterval(actualizarReloj, 1000);
        actualizarReloj();
    </script>
    """
    components.html(html_reloj, height=105)[cite: 7]

# --- GESTOR DE PERSISTENCIA ---
def cargar_estado():
    if not os.path.exists(FILE_DB):
        data_inicial = {
            "usuarios": {
                "admin1": {
                    "nombre": "Francesca Fellay",
                    "rol": "Admin",
                    "pin": "1234",
                    "permiso_editar": True,
                    "permiso_eliminar": True,
                    "avatar": ""
                },
                "gerente": {
                    "nombre": "Gerencia General",
                    "rol": "Admin",
                    "pin": "gerente",
                    "permiso_editar": True,
                    "permiso_eliminar": True,
                    "avatar": ""
                }
            },
            "datasets": {},
            "intervencion": [],
            "informes": []
        }
        guardar_estado(data_inicial)
        return data_inicial
    with open(FILE_DB, "r", encoding="utf-8") as f:
        try:
            data = json.load(f)
        except Exception:
            data = {"usuarios": {}, "datasets": {}, "intervencion": [], "informes": []}
    
    if "usuarios" not in data:
        data["usuarios"] = {}
    
    if "admin1" in data["usuarios"]:
        data["usuarios"]["admin1"]["nombre"] = "Francesca Fellay"
    else:
        data["usuarios"]["admin1"] = {
            "nombre": "Francesca Fellay",
            "rol": "Admin",
            "pin": "1234",
            "permiso_editar": True,
            "permiso_eliminar": True,
            "avatar": ""
        }
        
    data["usuarios"]["gerente"] = {
        "nombre": "Gerencia General",
        "rol": "Admin",
        "pin": "gerente",
        "permiso_editar": True,
        "permiso_eliminar": True,
        "avatar": ""
    }
    guardar_estado(data)
    return data[cite: 7]

def guardar_estado(data):
    with open(FILE_DB, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)[cite: 7]

db = cargar_estado()[cite: 7]

# --- CLIENTE GROQ IA ---
def obtener_api_key_groq():
    if "GROQ_API_KEY" in st.secrets:
        return st.secrets["GROQ_API_KEY"][cite: 6]
    return os.environ.get("GROQ_API_KEY", "")[cite: 6]

API_KEY_GROQ = obtener_api_key_groq()[cite: 6]

def obtener_cliente_ia():
    if not API_KEY_GROQ:
        return None
    try:
        return Groq(api_key=API_KEY_GROQ)[cite: 6]
    except Exception:
        return None

def ejecutar_chat_groq(client, prompt_sistema, prompt_usuario):
    modelos_candidatos = [
        "llama-3.3-70b-versatile",
        "llama-3.1-70b-versatile",
        "llama3-70b-8192",
        "mixtral-8x7b-32768",
        "gemma2-9b-it"
    ][cite: 6]
    try:
        lista_api = [m.id for m in client.models.list().data if "whisper" not in m.id.lower() and "guard" not in m.id.lower()][cite: 6]
    except Exception:
        lista_api = [][cite: 6]

    candidatos = [m for m in modelos_candidatos if m in lista_api] + lista_api + modelos_candidatos[cite: 6]
    candidatos = list(dict.fromkeys(candidatos))[cite: 6]

    ultimo_error = None
    for model_id in candidatos:
        try:
            resp = client.chat.completions.create(
                model=model_id,
                messages=[
                    {"role": "system", "content": prompt_sistema},
                    {"role": "user", "content": prompt_usuario}
                ],
                temperature=0.3
            )[cite: 6]
            return resp.choices[0].message.content[cite: 6]
        except Exception as e:
            ultimo_error = e
            continue
    raise Exception(f"No fue posible conectar con los modelos de Groq. Detalle: {ultimo_error}")[cite: 6]

# --- EXTRACCION DE TEXTO ---
def extraer_texto_archivo(ruta, extension):
    texto = ""
    try:
        if extension == "pdf":
            lector = pypdf.PdfReader(ruta)
            for pag in lector.pages:
                t = pag.extract_text()
                if t:
                    texto += t + "\n"
        elif extension == "docx":
            doc = Document(ruta)
            texto = "\n".join([p.text for p in doc.paragraphs])
        elif extension == "pptx":
            prs = Presentation(ruta)
            for diap in prs.slides:
                for forma in diap.shapes:
                    if hasattr(forma, "text"):
                        texto += forma.text + "\n"
        elif extension in ["xlsx", "xls"]:
            df_tmp = pd.read_excel(ruta)
            texto = f"Estadisticas:\n{df_tmp.describe(include='all').to_string()}\nPrimeras filas:\n{df_tmp.head(10).to_string()}"
    except Exception as e:
        texto = f"Error al extraer texto: {e}"
    return texto[:8000][cite: 7]

# --- CONTROL DE ACCESO (LOGIN) ---
if "autenticado" not in st.session_state:
    st.session_state.autenticado = False
    st.session_state.usuario_clave = None[cite: 7]

if not st.session_state.autenticado:
    st.markdown("<h1 style='text-align:center; color:#061e1b !important; margin-top:35px;'>🔐 Acceso a la Plataforma</h1>", unsafe_allow_html=True)
    st.markdown("<p style='text-align:center; color:#061e1b;'>Ingrese sus credenciales registradas</p>", unsafe_allow_html=True)
    c1, c2, c3 = st.columns([1, 1.4, 1])
    with c2:
        with st.form("login_form"):
            u_in = st.text_input("Usuario")
            p_in = st.text_input("PIN / Contrasena", type="password")
            if st.form_submit_button("Iniciar Sesion", use_container_width=True):
                if u_in in db["usuarios"] and db["usuarios"][u_in]["pin"] == p_in:
                    st.session_state.autenticado = True
                    st.session_state.usuario_clave = u_in
                    st.rerun()
                else:
                    st.error("Credenciales incorrectas")
    st.stop()[cite: 7]

# --- DATOS DEL USUARIO ACTUAL ---
usr_key = st.session_state.usuario_clave
usr = db["usuarios"].get(usr_key, {"nombre": "Invitado", "rol": "Usuario", "permiso_editar": False, "permiso_eliminar": False, "avatar": ""})
es_admin = (usr.get("rol") == "Admin") or (usr_key in ["admin1", "gerente"])[cite: 7]

# --- BARRA LATERAL ---
st.sidebar.markdown("### 👤 Mi Perfil")

avatar_path = usr.get("avatar", "")
if avatar_path and os.path.exists(avatar_path):
    st.sidebar.image(avatar_path, width=105)
else:
    st.sidebar.markdown("""<div style='width:80px; height:80px; border-radius:50%; background:#a2d2cc; border:2px solid #52948d; display:flex; align-items:center; justify-content:center; font-size:2.2rem; color:#061e1b; margin-bottom:10px;'>👤</div>""", unsafe_allow_html=True)

st.sidebar.markdown(f"**{usr.get('nombre')}**")
st.sidebar.markdown(f"<span class='badge-role'>{usr.get('rol')}</span>", unsafe_allow_html=True)

with st.sidebar.expander("📷 Cambiar foto de perfil"):
    nueva_foto = st.file_uploader("Subir imagen (JPG, PNG):", type=["jpg", "jpeg", "png"], key="upload_avatar")
    if nueva_foto is not None:
        if st.button("Guardar Foto"):
            ruta_avatar = os.path.join(DIR_AVATARS, f"{usr_key}_avatar.png")
            with open(ruta_avatar, "wb") as f_av:
                f_av.write(nueva_foto.getbuffer())
            db["usuarios"][usr_key]["avatar"] = ruta_avatar
            guardar_estado(db)
            st.success("Foto actualizada.")
            st.rerun()

st.sidebar.markdown("---")
if st.sidebar.button("🚪 Cerrar Sesion", use_container_width=True):
    st.session_state.autenticado = False
    st.session_state.usuario_clave = None
    st.rerun()[cite: 7]

# --- ENCABEZADO SUPERIOR CON RELOJ ---
col_head_title, col_head_clock = st.columns([1.2, 1])
with col_head_title:
    st.markdown("<h1 style='color:#061e1b !important; margin:0;'>⚡ Centro de Gestion & Analitica</h1>", unsafe_allow_html=True)
    st.markdown("<p style='color:#061e1b; margin-top:4px;'>Plataforma colaborativa multi-formato con procesamiento inteligente mediante Groq IA.</p>", unsafe_allow_html=True)
with col_head_clock:
    renderizar_reloj_chile()[cite: 7]

# --- NAVEGACION POR PESTANAS ---
titulos_pestanas = ["📊 Datasets & Archivos", "📂 Intervencion", "📄 Informes Compartidos"]
if es_admin:
    titulos_pestanas.append("👥 Gestion de Usuarios")

pestanas_principales = st.tabs(titulos_pestanas)[cite: 7]

# ==========================================
# 1. DATASETS EXCEL
# ==========================================
with pestanas_principales[0]:
    st.markdown("""<div class='modern-card'><h3 style='margin:0;'>📊 Repositorio de Datasets</h3><p style='margin:0; color:#061e1b;'>Suba, visualice y edite hojas de calculo en ventanas independientes.</p></div>""", unsafe_allow_html=True)
    
    with st.expander("➕ Cargar Nuevo Dataset Excel", expanded=True):
        c_tit, c_arc = st.columns([1, 1])
        with c_tit:
            t_data = st.text_input("Titulo del dataset:")
        with c_arc:
            f_data = st.file_uploader("Archivo Excel (.xlsx, .xls):", type=["xlsx", "xls"])
            
        if st.button("Guardar Dataset"):
            if not t_data.strip() or f_data is None:
                st.error("Complete el titulo y seleccione un archivo.")
            else:
                nom_arc = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{f_data.name}"
                ruta_dest = os.path.join(DIR_DATASETS, nom_arc)
                with open(ruta_dest, "wb") as f:
                    f.write(f_data.getbuffer())
                
                db["datasets"][t_data] = {
                    "titulo": t_data,
                    "autor": usr["nombre"],
                    "fecha": datetime.now().strftime("%Y-%m-%d %H:%M"),
                    "ruta": ruta_dest
                }
                guardar_estado(db)
                st.success("Dataset guardado con exito.")
                st.rerun()

    st.markdown("---")
    if not db["datasets"]:
        st.info("No hay datasets subidos actualmente.")
    else:
        titulos_ds = list(db["datasets"].keys())
        pestanas_ds = st.tabs(titulos_ds)
        
        for idx_ds, tab_ds in enumerate(pestanas_ds):
            t_act = titulos_ds[idx_ds]
            info_ds = db["datasets"][t_act]
            
            with tab_ds:
                st.markdown(f"### 📋 {info_ds['titulo']}")
                st.caption(f"Subido por: **{info_ds['autor']}** | Fecha: {info_ds['fecha']}")
                
                if os.path.exists(info_ds["ruta"]):
                    df_actual = pd.read_excel(info_ds["ruta"])
                    
                    if usr.get("permiso_editar") or es_admin:
                        st.markdown("**✏️ Editor de Datos en Vivo:**")
                        df_edit = st.data_editor(df_actual, key=f"d_edit_{t_act}", use_container_width=True)
                        if st.button("💾 Guardar Cambios en Excel", key=f"s_df_{t_act}"):
                            df_edit.to_excel(info_ds["ruta"], index=False)
                            st.success("Datos actualizados.")
                            st.rerun()
                    else:
                        st.dataframe(df_actual, use_container_width=True)
                    
                    if usr.get("permiso_eliminar") or es_admin:
                        st.markdown("---")
                        if st.button("🗑️ Eliminar Dataset", key=f"del_ds_{t_act}", type="secondary"):
                            if os.path.exists(info_ds["ruta"]):
                                os.remove(info_ds["ruta"])
                            del db["datasets"][t_act]
                            guardar_estado(db)
                            st.success("Dataset eliminado.")
                            st.rerun()
                else:
                    st.error("Archivo fisico no encontrado.")[cite: 7]

# ==========================================
# 2. INTERVENCION MULTI-FORMATO
# ==========================================
with pestanas_principales[1]:
    st.markdown("""<div class='modern-card'><h3 style='margin:0;'>📂 Modulo de Intervencion</h3><p style='margin:0; color:#061e1b;'>Soporte y resumenes automaticos con Groq para Documentos, Imagenes y Audios (M4A/MP3/WAV).</p></div>""", unsafe_allow_html=True)
    
    with st.form("form_intervencion"):
        tit_int = st.text_input("Titulo descriptivo del material:")
        archivos = st.file_uploader(
            "Cargar archivos:",
            type=["xlsx", "xls", "docx", "pdf", "pptx", "jpg", "jpeg", "png", "mp4", "m4a", "mp3", "wav"],
            accept_multiple_files=True
        )
        if st.form_submit_button("Subir y Procesar"):
            if not tit_int.strip() or not archivos:
                st.error("Complete el titulo y cargue al menos un archivo.")
            else:
                client = obtener_cliente_ia()
                for a in archivos:
                    ext = a.name.split(".")[-1].lower()
                    nom_dest = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{a.name}"
                    ruta_guardada = os.path.join(DIR_INTERVENCION, nom_dest)
                    
                    with open(ruta_guardada, "wb") as f:
                        f.write(a.getbuffer())
                    
                    resumen_txt = "Archivo guardado."
                    if client:
                        with st.spinner(f"Analizando {a.name} con Groq IA..."):
                            try:
                                if ext in ["m4a", "mp3", "wav"]:
                                    with open(ruta_guardada, "rb") as f_aud:
                                        audio_bytes = f_aud.read()
                                    audio_buffer = io.BytesIO(audio_bytes)
                                    audio_buffer.name = f"temp_audio.{ext}"
                                    
                                    transcripcion = client.audio.transcriptions.create(
                                        model=MODELO_WHISPER,
                                        file=audio_buffer,
                                        language="es"
                                    )[cite: 6]
                                    texto_audio = transcripcion.text[cite: 6]
                                    
                                    p_sys = "Eres un especialista senior en diagnostico e intervencion social/academica. Sintetiza con precision en espanol latinoamericano."
                                    p_user = f"A partir de la siguiente transcripcion de audio ({a.name}), sintetiza los puntos clave tratados, acuerdos y diagnostico para intervencion. Usa parrafos claros o vinetas Markdown estandar, nunca diagramas ASCII:\n\n{texto_audio}"
                                    resumen_txt = ejecutar_chat_groq(client, p_sys, p_user)
                                elif ext in ["pdf", "docx", "pptx", "xlsx", "xls"]:
                                    t_doc = extraer_texto_archivo(ruta_guardada, ext)
                                    p_sys = "Eres un especialista senior en analisis documental e intervencion."
                                    p_user = f"Elabora un resumen y diagnostico clave del siguiente documento ({a.name}):\n\n{t_doc}\n\nUsa vinetas Markdown estandar con conceptos clave en negrita."
                                    resumen_txt = ejecutar_chat_groq(client, p_sys, p_user)
                                elif ext in ["jpg", "jpeg", "png"]:
                                    resumen_txt = f"Imagen registrada ({a.name}). Visualizacion disponible en la pestana multimedia."
                                elif ext == "mp4":
                                    resumen_txt = f"Video registrado ({a.name}). Reproduccion multimedia disponible."
                            except Exception as e:
                                resumen_txt = f"Archivo guardado. Diagnostico no generado: {e}"
                    
                    db["intervencion"].append({
                        "titulo": tit_int,
                        "nombre_original": a.name,
                        "tipo": ext,
                        "autor": usr["nombre"],
                        "fecha": datetime.now().strftime("%Y-%m-%d %H:%M"),
                        "ruta": ruta_guardada,
                        "resumen": resumen_txt
                    })
                guardar_estado(db)
                st.success("Materiales integrados correctamente.")
                st.rerun()

    st.markdown("---")
    st.subheader("📚 Materiales Guardados")
    
    if not db["intervencion"]:
        st.info("No hay registros en intervencion.")
    else:
        for idx, item in enumerate(db["intervencion"]):
            titulo_item = item.get("titulo") or item.get("titulo_registro") or "Sin Titulo"
            nombre_arc = item.get("nombre_original") or item.get("filename") or "Archivo"
            tipo_arc = item.get("tipo", "").lower()
            autor_arc = item.get("autor", "Desconocido")
            fecha_arc = item.get("fecha", "")
            resumen_arc = item.get("resumen", "Sin resumen disponible.")
            ruta_arc = item.get("ruta", "")
            
            with st.container():
                col_head, col_del_btn = st.columns([5, 1])
                with col_head:
                    st.markdown(f"### 📁 {titulo_item} — <span class='highlight-tag'>{nombre_arc}</span>", unsafe_allow_html=True)
                    st.caption(f"Autor: **{autor_arc}** | Fecha: {fecha_arc} | Formato: **{tipo_arc.upper()}**")
                with col_del_btn:
                    if usr.get("permiso_eliminar") or es_admin:
                        if st.button("🗑️ Borrar", key=f"del_int_{idx}", type="secondary"):
                            if ruta_arc and os.path.exists(ruta_arc):
                                os.remove(ruta_arc)
                            db["intervencion"].pop(idx)
                            guardar_estado(db)
                            st.success("Archivo eliminado.")
                            st.rerun()
                
                t_vis, t_res = st.tabs(["👁️ Multimedia / Descarga", "📝 Diagnostico y Resumen"])
                with t_vis:
                    if ruta_arc and os.path.exists(ruta_arc):
                        if tipo_arc == "mp4":
                            st.video(ruta_arc)
                        elif tipo_arc in ["m4a", "mp3", "wav"]:
                            st.audio(ruta_arc)
                        elif tipo_arc in ["jpg", "jpeg", "png"]:
                            st.image(ruta_arc, use_container_width=True)
                        else:
                            with open(ruta_arc, "rb") as fl:
                                st.download_button("📥 Descargar Archivo", data=fl.read(), file_name=nombre_arc, key=f"dl_a_{idx}")
                    else:
                        st.error("Archivo fisico no encontrado.")
                with t_res:
                    st.markdown(resumen_arc)
                st.markdown("---")[cite: 7]

# ==========================================
# 3. INFORMES COMPARTIDOS
# ==========================================
with pestanas_principales[2]:
    st.markdown("""<div class='modern-card'><h3 style='margin:0;'>📄 Informes Compartidos</h3><p style='margin:0; color:#061e1b;'>Generacion y repositorio colaborativo con exportacion en formato Carta.</p></div>""", unsafe_allow_html=True)
    
    with st.expander("🤖 Redactar Nuevo Informe con IA (Groq)", expanded=False):
        nom_i = st.text_input("Titulo del Informe:")
        enf_i = st.selectbox("Enfoque:", ["Resumen Ejecutivo", "Diagnostico Tecnico", "Evaluacion Estrategica"])
        ins_i = st.text_area("Instrucciones complementarias:")
        
        if st.button("🚀 Generar Informe"):
            client = obtener_cliente_ia()
            if not nom_i.strip():
                st.error("Ingrese un titulo para el informe.")
            elif not client:
                st.error("API Key de Groq (GROQ_API_KEY) no configurada en Secrets.")
            else:
                with st.spinner("Redactando informe consolidado con Groq IA..."):
                    resumenes_int = "\n".join([f"- {x.get('nombre_original', 'Archivo')}: {x.get('resumen', '')}" for x in db["intervencion"]])
                    
                    p_sys = "Actua como especialista analitico senior. Genera informes estructurados de excelencia en espanol latinoamericano."
                    p_user = f"""Genera un informe estructurado profesional con enfoque '{enf_i}'.
Titulo: {nom_i}
Autor solicitante: {usr['nombre']}
Instrucciones: {ins_i}
Materiales analizados: {resumenes_int if resumenes_int else 'Sin materiales adjuntos.'}

ESTRUCTURA REQUERIDA:
1. Diagnostico y Vision Global
2. Hallazgos Analiticos Relevantes (Si incluyes tablas comparativas o sintesis, hazlo OBLIGATORIAMENTE usando sintaxis estandar de tablas Markdown con '|' y encabezados separados por '|---|---|', NUNCA uses arte ASCII ni bloques de codigo con ``` para tablas).
3. Conclusiones y Plan de Accion

REGLA OBLIGATORIA DE FORMATO:
- Las tablas deben ser tablas Markdown renderizables reales.
- Todo el texto analitico debe fluir con negritas y vinetas claras."""
                    try:
                        resp_txt = ejecutar_chat_groq(client, p_sys, p_user)
                        db["informes"].append({
                            "titulo": nom_i,
                            "autor": usr["nombre"],
                            "fecha": datetime.now().strftime("%Y-%m-%d %H:%M"),
                            "contenido": resp_txt
                        })
                        guardar_estado(db)
                        st.success("Informe publicado.")
                        st.rerun()
                    except Exception as err:
                        st.error(f"Error al generar informe: {err}")

    st.markdown("---")
    st.subheader("📚 Repositorio de Informes")
    
    if not db["informes"]:
        st.info("No hay informes registrados.")
    else:
        for idx_inf, inf in enumerate(db["informes"]):
            titulo_inf = inf.get("titulo") or inf.get("titulo_informe") or "Sin Titulo"
            autor_inf = inf.get("autor", "Desconocido")
            fecha_inf = inf.get("fecha", "")
            contenido_inf = inf.get("contenido", "")
            
            with st.container():
                c_inf_t, c_inf_del = st.columns([5, 1])
                with c_inf_t:
                    st.markdown(f"### 📄 [{autor_inf}] - {titulo_inf}")
                    st.caption(f"Generado el: {fecha_inf}")
                with c_inf_del:
                    if usr.get("permiso_eliminar") or es_admin:
                        if st.button("🗑️ Borrar", key=f"del_inf_{idx_inf}", type="secondary"):
                            db["informes"].pop(idx_inf)
                            guardar_estado(db)
                            st.success("Informe eliminado.")
                            st.rerun()
                
                st.markdown(contenido_inf)
                
                cd1, cd2 = st.columns([1, 1])
                with cd1:
                    df_out = pd.DataFrame([{"Autor": autor_inf, "Titulo": titulo_inf, "Fecha": fecha_inf, "Contenido": contenido_inf}])
                    buf = io.BytesIO()
                    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
                        df_out.to_excel(writer, index=False)
                    st.download_button(
                        label="📥 Descargar Excel (.xlsx)",
                        data=buf.getvalue(),
                        file_name=f"[{autor_inf}] - {titulo_inf}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        key=f"xls_dl_{idx_inf}"
                    )
                with cd2:
                    html_carta = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{titulo_inf}</title>
    <style>
        @page {{ size: letter portrait; margin: 25mm; }}
        body {{ font-family: 'Segoe UI', Arial, sans-serif; color: #061e1b; padding: 20px; background-color: #89c7c0; }}
        .header {{ border-bottom: 2px solid #52948d; padding-bottom: 8px; margin-bottom: 20px; }}
        .title {{ font-size: 20pt; font-weight: bold; color: #061e1b; }}
        .meta {{ font-size: 10pt; color: #061e1b; margin-top: 4px; }}
        .body {{ font-size: 11pt; line-height: 1.6; white-space: pre-wrap; }}
        table {{ width: 100%; border-collapse: collapse; margin: 15px 0; }}
        th, td {{ border: 1px solid #52948d; padding: 8px 12px; text-align: left; }}
        th {{ background-color: #63a59e; }}
    </style>
</head>
<body>
    <div class="header">
        <div class="title">{titulo_inf}</div>
        <div class="meta">Autor: {autor_inf} | Fecha: {fecha_inf}</div>
    </div>
    <div class="body">{contenido_inf}</div>
</body>
</html>"""
                    st.download_button(
                        label="🖼️ Descargar Formato Carta (HTML / PDF)",
                        data=html_carta,
                        file_name=f"[{autor_inf}] - {titulo_inf}_Carta.html",
                        mime="text/html",
                        key=f"doc_dl_{idx_inf}"
                    )
                st.markdown("---")[cite: 7]

# ==========================================
# 4. GESTION DE USUARIOS (SOLO ADMIN)
# ==========================================
if es_admin:
    with pestanas_principales[3]:
        st.markdown("""<div class='modern-card'><h3 style='margin:0;'>👥 Gestion de Usuarios y Permisos</h3><p style='margin:0; color:#061e1b;'>Control centralizado de cuentas accesible exclusivamente por Administradores.</p></div>""", unsafe_allow_html=True)
        
        with st.expander("➕ Registrar Nuevo Usuario", expanded=True):
            cu1, cu2, cu3 = st.columns([1, 1, 1])
            with cu1:
                n_u = st.text_input("Usuario (Login):")
                n_nom = st.text_input("Nombre Completo:")
            with cu2:
                n_pin = st.text_input("PIN / Clave:", type="password")
                n_rol = st.selectbox("Rol:", ["Usuario", "Analista", "Especialista", "Gerencia", "Admin"])
            with cu3:
                st.markdown("**Permisos Iniciales:**")
                p_e = st.checkbox("Permiso para Editar Datasets")
                p_d = st.checkbox("Permiso para Eliminar Datasets / Archivos")
                
            if st.button("Crear Usuario"):
                if not n_u or not n_pin or not n_nom:
                    st.error("Todos los campos son obligatorios.")
                elif n_u in db["usuarios"]:
                    st.error("El nombre de usuario ya existe.")
                else:
                    db["usuarios"][n_u] = {
                        "nombre": n_nom,
                        "rol": n_rol,
                        "pin": n_pin,
                        "permiso_editar": p_e,
                        "permiso_eliminar": p_d,
                        "avatar": ""
                    }
                    guardar_estado(db)
                    st.success(f"Usuario '{n_u}' registrado correctamente.")
                    st.rerun()

        st.markdown("---")
        st.subheader("📜 Cuentas Registradas")
        
        usuarios_visibles = [k for k in list(db["usuarios"].keys()) if k != "gerente"]
        
        for u_k in usuarios_visibles:
            u_d = db["usuarios"][u_k]
            with st.container():
                col_u_inf, col_u_e, col_u_d, col_u_del = st.columns([2, 1, 1, 1])
                with col_u_inf:
                    st.markdown(f"**{u_d.get('nombre', '')}** (<span class='highlight-tag'>{u_k}</span>) — Rol: *{u_d.get('rol',
